#!/usr/bin/env python3
"""
Professional PowerPoint Creator - Main Generation Engine
Transforms user prompts into comprehensive, professional presentations
"""

import json
import sys
import re
from typing import Dict, List, Tuple, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

class ProfessionalPPTCreator:
    """Professional PowerPoint presentation generator"""
    
    def __init__(self):
        self.color_schemes = {
            'corporate': {
                'primary': RGBColor(0, 51, 102),    # Deep blue
                'secondary': RGBColor(0, 102, 204), # Bright blue
                'accent': RGBColor(255, 165, 0),    # Orange
                'text': RGBColor(51, 51, 51),       # Dark gray
                'background': RGBColor(255, 255, 255) # White
            },
            'modern': {
                'primary': RGBColor(45, 45, 45),    # Charcoal
                'secondary': RGBColor(100, 200, 200), # Teal
                'accent': RGBColor(255, 107, 107),  # Coral
                'text': RGBColor(70, 70, 70),       # Medium gray
                'background': RGBColor(248, 249, 250) # Light gray
            },
            'creative': {
                'primary': RGBColor(88, 24, 69),    # Deep purple
                'secondary': RGBColor(144, 224, 239), # Light blue
                'accent': RGBColor(255, 195, 0),    # Golden yellow
                'text': RGBColor(40, 40, 40),       # Dark gray
                'background': RGBColor(255, 255, 255) # White
            }
        }
        
    def analyze_prompt(self, prompt: str) -> Dict:
        """Analyze user prompt to extract topic, audience, and structure requirements"""
        # Detect presentation type
        business_keywords = ['business', 'board', 'executive', 'strategy', 'financial', 'quarterly', 'revenue']
        educational_keywords = ['training', 'course', 'lesson', 'workshop', 'seminar', 'tutorial']
        sales_keywords = ['pitch', 'proposal', 'client', 'product', 'solution', 'demo']
        
        prompt_lower = prompt.lower()
        
        if any(word in prompt_lower for word in business_keywords):
            pres_type = 'business'
        elif any(word in prompt_lower for word in educational_keywords):
            pres_type = 'educational'
        elif any(word in prompt_lower for word in sales_keywords):
            pres_type = 'sales'
        else:
            pres_type = 'general'
        
        # Extract topic (simplified - in real implementation would use NLP)
        topic = self._extract_main_topic(prompt)
        
        # Determine slide count
        slide_count = self._estimate_slide_count(prompt, pres_type)
        
        return {
            'topic': topic,
            'type': pres_type,
            'slide_count': slide_count,
            'color_scheme': 'corporate' if pres_type == 'business' else 'modern',
            'audience': self._detect_audience(prompt)
        }
    
    def _extract_main_topic(self, prompt: str) -> str:
        """Extract main topic from prompt"""
        # Remove common words and find key topic
        stop_words = {'create', 'make', 'presentation', 'about', 'on', 'for', 'the', 'a', 'an'}
        words = [word.strip('.,!?') for word in prompt.split()]
        topic_words = [word for word in words if word.lower() not in stop_words and len(word) > 2]
        return ' '.join(topic_words[:3])  # Take first 3 meaningful words
    
    def _estimate_slide_count(self, prompt: str, pres_type: str) -> int:
        """Estimate appropriate number of slides based on prompt complexity"""
        word_count = len(prompt.split())
        
        if word_count < 10:
            return 5  # Brief prompt = concise presentation
        elif word_count < 25:
            return 8  # Medium prompt = standard presentation  
        else:
            return 12  # Detailed prompt = comprehensive presentation
    
    def _detect_audience(self, prompt: str) -> str:
        """Detect intended audience from prompt context"""
        if any(word in prompt.lower() for word in ['board', 'executive', 'c-level', 'senior']):
            return 'executives'
        elif any(word in prompt.lower() for word in ['team', 'colleagues', 'staff', 'employees']):
            return 'internal'
        elif any(word in prompt.lower() for word in ['client', 'customer', 'prospect', 'buyer']):
            return 'external'
        else:
            return 'general'
    
    def generate_presentation(self, prompt: str, output_path: str = None) -> str:
        """Generate complete professional presentation from prompt"""
        analysis = self.analyze_prompt(prompt)
        
        # Create presentation
        prs = Presentation()
        
        # Generate slides based on analysis
        slides_content = self._generate_slide_content(analysis, prompt)
        
        # Create slides
        for i, slide_content in enumerate(slides_content):
            if i == 0:
                slide = self._create_title_slide(prs, slide_content)
            elif 'agenda' in slide_content.get('type', '').lower():
                slide = self._create_agenda_slide(prs, slide_content)
            elif 'data' in slide_content.get('type', '').lower():
                slide = self._create_data_slide(prs, slide_content)
            else:
                slide = self._create_content_slide(prs, slide_content)
            
            # Apply color scheme
            self._apply_color_scheme(slide, analysis['color_scheme'])
        
        # Save presentation
        if not output_path:
            safe_topic = re.sub(r'[^\w\s-]', '', analysis['topic']).strip()
            safe_topic = re.sub(r'[-\s]+', '_', safe_topic)
            output_path = f"{safe_topic}_presentation.pptx"
        
        prs.save(output_path)
        return output_path
    
    def _generate_slide_content(self, analysis: Dict, original_prompt: str) -> List[Dict]:
        """Generate content for each slide based on analysis"""
        slides = []
        
        # Title slide
        slides.append({
            'type': 'title',
            'title': analysis['topic'],
            'subtitle': f"Professional Presentation • {analysis['audience'].title()} Focus"
        })
        
        # Agenda/Overview slide
        if analysis['slide_count'] > 5:
            slides.append({
                'type': 'agenda',
                'title': 'Agenda',
                'bullets': self._generate_agenda_items(analysis)
            })
        
        # Main content slides
        content_slides = self._generate_main_content(analysis, original_prompt)
        slides.extend(content_slides)
        
        # Conclusion slide
        slides.append({
            'type': 'conclusion',
            'title': 'Next Steps',
            'bullets': self._generate_next_steps(analysis)
        })
        
        return slides
    
    def _generate_agenda_items(self, analysis: Dict) -> List[str]:
        """Generate agenda items based on presentation type"""
        if analysis['type'] == 'business':
            return [
                "Executive Summary",
                "Current Situation Analysis",
                "Strategic Recommendations", 
                "Implementation Timeline",
                "Expected Outcomes",
                "Q&A Discussion"
            ]
        elif analysis['type'] == 'sales':
            return [
                "Challenge Overview",
                "Solution Framework",
                "Benefits & ROI",
                "Implementation Process",
                "Success Stories",
                "Next Steps"
            ]
        else:
            return [
                "Introduction",
                "Key Concepts", 
                "Detailed Analysis",
                "Practical Applications",
                "Summary",
                "Questions"
            ]
    
    def _generate_main_content(self, analysis: Dict, prompt: str) -> List[Dict]:
        """Generate main content slides"""
        slides = []
        target_slides = analysis['slide_count'] - 3  # Minus title, agenda, conclusion
        
        for i in range(target_slides):
            slides.append({
                'type': 'content',
                'title': f"Key Point {i+1}",
                'bullets': [
                    f"Supporting detail about {analysis['topic']}",
                    f"Evidence and examples relevant to {analysis['audience']}",
                    f"Implications and impact analysis",
                    f"Action items and recommendations"
                ]
            })
        
        return slides
    
    def _generate_next_steps(self, analysis: Dict) -> List[str]:
        """Generate next steps based on audience and type"""
        if analysis['audience'] == 'executives':
            return [
                "Review strategic recommendations",
                "Approve resource allocation",
                "Set implementation timeline",
                "Schedule follow-up meeting"
            ]
        elif analysis['type'] == 'sales':
            return [
                "Schedule detailed solution demo",
                "Provide custom proposal",
                "Connect with implementation team",
                "Begin pilot program discussion"
            ]
        else:
            return [
                "Apply key concepts learned",
                "Schedule follow-up session", 
                "Share feedback and questions",
                "Continue skill development"
            ]
    
    def _create_title_slide(self, prs: Presentation, content: Dict):
        """Create professional title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content['title']
        subtitle.text = content['subtitle']
        
        # Style title
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Style subtitle  
        subtitle_para = subtitle.text_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_agenda_slide(self, prs: Presentation, content: Dict):
        """Create agenda slide with bullet points"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = content['title']
        
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        for bullet in content['bullets']:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
        
        return slide
    
    def _create_content_slide(self, prs: Presentation, content: Dict):
        """Create standard content slide"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title  
        title.text = content['title']
        
        if 'bullets' in content:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for bullet in content['bullets']:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
        
        return slide
    
    def _create_data_slide(self, prs: Presentation, content: Dict):
        """Create slide optimized for data presentation"""
        # Implementation for charts and data visualization
        return self._create_content_slide(prs, content)
    
    def _apply_color_scheme(self, slide, scheme_name: str):
        """Apply consistent color scheme to slide"""
        scheme = self.color_schemes.get(scheme_name, self.color_schemes['corporate'])
        
        # Apply to title if exists
        if hasattr(slide, 'shapes') and slide.shapes.title:
            title_para = slide.shapes.title.text_frame.paragraphs[0]
            title_para.font.color.rgb = scheme['primary']

def main():
    """Command line interface for presentation generation"""
    if len(sys.argv) < 2:
        print("Usage: python ppt_creator.py 'Your presentation topic/prompt'")
        print("Example: python ppt_creator.py 'Create a business presentation about renewable energy for board meeting'")
        return
    
    prompt = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    creator = ProfessionalPPTCreator()
    
    try:
        result_path = creator.generate_presentation(prompt, output_file)
        print(f"✅ Professional presentation created: {result_path}")
        
        # Output metadata for Claude
        analysis = creator.analyze_prompt(prompt)
        print(f"📊 Generated {analysis['slide_count']} slides for {analysis['audience']} audience")
        print(f"🎨 Applied {analysis['color_scheme']} color scheme")
        print(f"📋 Presentation type: {analysis['type']}")
        
    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()