#!/usr/bin/env python3
"""
Enhanced Professional PowerPoint Creator - Advanced Generation Engine
Provides full customization control with advanced themes, animations, multimedia, and editing
"""

import json
import sys
import re
from typing import Dict, List, Tuple, Optional, Any
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE
from pptx.oxml.xmlchemy import OxmlElement
import io
import os

class AdvancedTheme:
    """Advanced theme configuration with full customization"""

    def __init__(self, name: str, config: Dict):
        self.name = name
        self.primary_color = config.get('primary', RGBColor(0, 51, 102))
        self.secondary_color = config.get('secondary', RGBColor(0, 102, 204))
        self.accent_color = config.get('accent', RGBColor(255, 165, 0))
        self.text_color = config.get('text', RGBColor(51, 51, 51))
        self.background_color = config.get('background', RGBColor(255, 255, 255))
        self.title_font = config.get('title_font', 'Calibri')
        self.body_font = config.get('body_font', 'Calibri')
        self.title_size = config.get('title_size', 44)
        self.body_size = config.get('body_size', 18)
        self.has_gradient = config.get('gradient', False)
        self.gradient_angle = config.get('gradient_angle', 90)

class EnhancedPPTCreator:
    """Enhanced PowerPoint creator with advanced customization"""

    def __init__(self):
        self.themes = self._initialize_themes()
        self.transition_types = ['none', 'fade', 'push', 'wipe', 'split', 'reveal', 'random_bars', 'shape', 'uncover', 'cover', 'flash', 'dissolve']
        self.animation_types = ['appear', 'fade', 'fly_in', 'float_in', 'split', 'wipe', 'wheel', 'random_bars', 'grow_and_turn', 'zoom', 'swivel', 'bounce']

    def _initialize_themes(self) -> Dict[str, AdvancedTheme]:
        """Initialize 10+ professional themes with full customization"""
        themes = {
            'corporate_blue': AdvancedTheme('Corporate Blue', {
                'primary': RGBColor(0, 51, 102),
                'secondary': RGBColor(0, 102, 204),
                'accent': RGBColor(255, 165, 0),
                'text': RGBColor(51, 51, 51),
                'background': RGBColor(255, 255, 255),
                'title_font': 'Calibri',
                'body_font': 'Calibri',
                'title_size': 44,
                'body_size': 18
            }),
            'modern_minimal': AdvancedTheme('Modern Minimal', {
                'primary': RGBColor(45, 45, 45),
                'secondary': RGBColor(100, 200, 200),
                'accent': RGBColor(255, 107, 107),
                'text': RGBColor(70, 70, 70),
                'background': RGBColor(248, 249, 250),
                'title_font': 'Arial',
                'body_font': 'Arial',
                'title_size': 40,
                'body_size': 16
            }),
            'creative_bold': AdvancedTheme('Creative Bold', {
                'primary': RGBColor(88, 24, 69),
                'secondary': RGBColor(144, 224, 239),
                'accent': RGBColor(255, 195, 0),
                'text': RGBColor(40, 40, 40),
                'background': RGBColor(255, 255, 255),
                'title_font': 'Impact',
                'body_font': 'Calibri',
                'title_size': 48,
                'body_size': 18
            }),
            'tech_startup': AdvancedTheme('Tech Startup', {
                'primary': RGBColor(16, 185, 129),
                'secondary': RGBColor(59, 130, 246),
                'accent': RGBColor(249, 115, 22),
                'text': RGBColor(17, 24, 39),
                'background': RGBColor(255, 255, 255),
                'title_font': 'Montserrat',
                'body_font': 'Open Sans',
                'title_size': 42,
                'body_size': 18,
                'gradient': True
            }),
            'elegant_dark': AdvancedTheme('Elegant Dark', {
                'primary': RGBColor(236, 72, 153),
                'secondary': RGBColor(147, 51, 234),
                'accent': RGBColor(245, 158, 11),
                'text': RGBColor(229, 231, 235),
                'background': RGBColor(17, 24, 39),
                'title_font': 'Georgia',
                'body_font': 'Georgia',
                'title_size': 44,
                'body_size': 18
            }),
            'finance_professional': AdvancedTheme('Finance Professional', {
                'primary': RGBColor(21, 94, 117),
                'secondary': RGBColor(82, 183, 136),
                'accent': RGBColor(253, 203, 110),
                'text': RGBColor(45, 55, 72),
                'background': RGBColor(255, 255, 255),
                'title_font': 'Times New Roman',
                'body_font': 'Arial',
                'title_size': 40,
                'body_size': 16
            }),
            'healthcare_calm': AdvancedTheme('Healthcare Calm', {
                'primary': RGBColor(59, 130, 246),
                'secondary': RGBColor(139, 92, 246),
                'accent': RGBColor(16, 185, 129),
                'text': RGBColor(55, 65, 81),
                'background': RGBColor(249, 250, 251),
                'title_font': 'Verdana',
                'body_font': 'Verdana',
                'title_size': 38,
                'body_size': 16
            }),
            'education_bright': AdvancedTheme('Education Bright', {
                'primary': RGBColor(239, 68, 68),
                'secondary': RGBColor(59, 130, 246),
                'accent': RGBColor(245, 158, 11),
                'text': RGBColor(31, 41, 55),
                'background': RGBColor(254, 252, 232),
                'title_font': 'Comic Sans MS',
                'body_font': 'Trebuchet MS',
                'title_size': 44,
                'body_size': 20
            }),
            'luxury_gold': AdvancedTheme('Luxury Gold', {
                'primary': RGBColor(180, 83, 9),
                'secondary': RGBColor(217, 119, 6),
                'accent': RGBColor(252, 211, 77),
                'text': RGBColor(120, 53, 15),
                'background': RGBColor(254, 249, 231),
                'title_font': 'Garamond',
                'body_font': 'Garamond',
                'title_size': 46,
                'body_size': 18
            }),
            'nature_organic': AdvancedTheme('Nature Organic', {
                'primary': RGBColor(34, 197, 94),
                'secondary': RGBColor(132, 204, 22),
                'accent': RGBColor(251, 191, 36),
                'text': RGBColor(22, 101, 52),
                'background': RGBColor(247, 254, 231),
                'title_font': 'Century Gothic',
                'body_font': 'Century Gothic',
                'title_size': 42,
                'body_size': 18
            }),
            'monochrome_professional': AdvancedTheme('Monochrome Professional', {
                'primary': RGBColor(0, 0, 0),
                'secondary': RGBColor(75, 85, 99),
                'accent': RGBColor(156, 163, 175),
                'text': RGBColor(31, 41, 55),
                'background': RGBColor(255, 255, 255),
                'title_font': 'Arial',
                'body_font': 'Arial',
                'title_size': 40,
                'body_size': 16
            }),
            'sunset_vibrant': AdvancedTheme('Sunset Vibrant', {
                'primary': RGBColor(244, 63, 94),
                'secondary': RGBColor(251, 146, 60),
                'accent': RGBColor(251, 191, 36),
                'text': RGBColor(127, 29, 29),
                'background': RGBColor(255, 251, 235),
                'title_font': 'Tahoma',
                'body_font': 'Tahoma',
                'title_size': 44,
                'body_size': 18,
                'gradient': True,
                'gradient_angle': 45
            })
        }
        return themes

    def create_presentation(self, config: Dict) -> str:
        """
        Create presentation with full manual control

        Args:
            config: Complete configuration dictionary with:
                - title: Presentation title
                - subtitle: Presentation subtitle (optional)
                - theme: Theme name or custom theme config
                - slides: List of slide configurations
                - transitions: Global or per-slide transitions
                - animations: Global or per-slide animations
                - output_path: Where to save the file
        """
        # Initialize presentation
        prs = Presentation()

        # Get theme
        theme = self._get_theme(config.get('theme', 'corporate_blue'))

        # Get global settings
        global_transition = config.get('global_transition', 'none')
        global_animation = config.get('global_animation', 'none')

        # Create slides
        slides_config = config.get('slides', [])

        for i, slide_config in enumerate(slides_config):
            slide = self._create_slide(prs, slide_config, theme)

            # Apply transitions
            transition = slide_config.get('transition', global_transition)
            if transition != 'none':
                self._add_transition(slide, transition, slide_config.get('transition_speed', 'medium'))

            # Apply animations would be handled here (note: python-pptx has limited animation support)
            # For full animation control, you'd need to manipulate the XML directly

        # Save presentation
        output_path = config.get('output_path', 'presentation.pptx')
        prs.save(output_path)

        return output_path

    def edit_presentation(self, file_path: str, modifications: Dict) -> str:
        """
        Edit an existing presentation

        Args:
            file_path: Path to existing presentation
            modifications: Dictionary of modifications to apply:
                - change_theme: New theme to apply
                - update_slides: Dict of slide_index: new_content
                - reorder_slides: List of new slide order
                - add_slides: List of new slide configs to add
                - delete_slides: List of slide indices to delete
        """
        # Load existing presentation
        prs = Presentation(file_path)

        # Change theme if requested
        if 'change_theme' in modifications:
            new_theme = self._get_theme(modifications['change_theme'])
            self._apply_theme_to_presentation(prs, new_theme)

        # Update specific slides
        if 'update_slides' in modifications:
            for slide_idx, new_content in modifications['update_slides'].items():
                if slide_idx < len(prs.slides):
                    slide = prs.slides[slide_idx]
                    self._update_slide_content(slide, new_content)

        # Add new slides
        if 'add_slides' in modifications:
            theme = self._get_theme(modifications.get('theme', 'corporate_blue'))
            for slide_config in modifications['add_slides']:
                self._create_slide(prs, slide_config, theme)

        # Delete slides (in reverse order to maintain indices)
        if 'delete_slides' in modifications:
            for slide_idx in sorted(modifications['delete_slides'], reverse=True):
                if slide_idx < len(prs.slides):
                    rId = prs.slides._sldIdLst[slide_idx].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[slide_idx]

        # Reorder slides
        if 'reorder_slides' in modifications:
            new_order = modifications['reorder_slides']
            self._reorder_slides(prs, new_order)

        # Save with new name or overwrite
        output_path = modifications.get('output_path', file_path.replace('.pptx', '_edited.pptx'))
        prs.save(output_path)

        return output_path

    def _get_theme(self, theme_input: Any) -> AdvancedTheme:
        """Get theme from name or custom configuration"""
        if isinstance(theme_input, str):
            return self.themes.get(theme_input, self.themes['corporate_blue'])
        elif isinstance(theme_input, dict):
            return AdvancedTheme('custom', theme_input)
        return self.themes['corporate_blue']

    def _create_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create a slide with full customization"""
        slide_type = config.get('type', 'content')

        if slide_type == 'title':
            return self._create_title_slide(prs, config, theme)
        elif slide_type == 'section':
            return self._create_section_slide(prs, config, theme)
        elif slide_type == 'content':
            return self._create_content_slide(prs, config, theme)
        elif slide_type == 'two_column':
            return self._create_two_column_slide(prs, config, theme)
        elif slide_type == 'image':
            return self._create_image_slide(prs, config, theme)
        elif slide_type == 'comparison':
            return self._create_comparison_slide(prs, config, theme)
        elif slide_type == 'timeline':
            return self._create_timeline_slide(prs, config, theme)
        elif slide_type == 'blank':
            return self._create_blank_slide(prs, config, theme)
        else:
            return self._create_content_slide(prs, config, theme)

    def _create_title_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create title slide with theme"""
        slide_layout = prs.slide_layouts[6]  # Blank layout for full control
        slide = prs.slides.add_slide(slide_layout)

        # Apply background
        self._apply_background(slide, theme, config.get('background'))

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(9), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = config.get('title', 'Presentation Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.name = theme.title_font
        title_para.font.size = Pt(theme.title_size)
        title_para.font.color.rgb = theme.primary_color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if 'subtitle' in config:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5),
                Inches(9), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = config['subtitle']
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.name = theme.body_font
            subtitle_para.font.size = Pt(theme.body_size)
            subtitle_para.font.color.rgb = theme.text_color
            subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    def _create_section_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create section divider slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Full color background
        self._apply_solid_background(slide, theme.primary_color)

        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(8), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = config.get('title', 'Section')
        title_para = title_frame.paragraphs[0]
        title_para.font.name = theme.title_font
        title_para.font.size = Pt(theme.title_size + 4)
        title_para.font.color.rgb = theme.background_color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        return slide

    def _create_content_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create standard content slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._apply_background(slide, theme, config.get('background'))

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = config.get('title', 'Slide Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.name = theme.title_font
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = theme.primary_color
        title_para.font.bold = True

        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Add bullets
        bullets = config.get('bullets', [])
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = bullet
            p.font.name = theme.body_font
            p.font.size = Pt(theme.body_size)
            p.font.color.rgb = theme.text_color
            p.level = 0
            p.space_after = Pt(12)

        return slide

    def _create_two_column_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create two-column layout slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._apply_background(slide, theme, config.get('background'))

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = config.get('title', 'Two Column Slide')
        title_para = title_frame.paragraphs[0]
        title_para.font.name = theme.title_font
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = theme.primary_color
        title_para.font.bold = True

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(4.25), Inches(4.5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        left_bullets = config.get('left_content', [])
        for i, bullet in enumerate(left_bullets):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = bullet
            p.font.name = theme.body_font
            p.font.size = Pt(theme.body_size - 2)
            p.font.color.rgb = theme.text_color

        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(5.25), Inches(1.5),
            Inches(4.25), Inches(4.5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        right_bullets = config.get('right_content', [])
        for i, bullet in enumerate(right_bullets):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = bullet
            p.font.name = theme.body_font
            p.font.size = Pt(theme.body_size - 2)
            p.font.color.rgb = theme.text_color

        return slide

    def _create_comparison_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create comparison slide with vs layout"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._apply_background(slide, theme, config.get('background'))

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = config.get('title', 'Comparison')
        title_para = title_frame.paragraphs[0]
        title_para.font.name = theme.title_font
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = theme.primary_color
        title_para.font.bold = True

        # Left box
        left_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.7),
            Inches(4), Inches(4)
        )
        left_shape.fill.solid()
        left_shape.fill.fore_color.rgb = theme.secondary_color
        left_text = left_shape.text_frame
        left_text.text = config.get('left_title', 'Option A')

        # VS text
        vs_box = slide.shapes.add_textbox(
            Inches(4.5), Inches(3.5),
            Inches(1), Inches(0.5)
        )
        vs_frame = vs_box.text_frame
        vs_frame.text = "VS"
        vs_para = vs_frame.paragraphs[0]
        vs_para.font.bold = True
        vs_para.font.size = Pt(24)
        vs_para.alignment = PP_ALIGN.CENTER

        # Right box
        right_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.5), Inches(1.7),
            Inches(4), Inches(4)
        )
        right_shape.fill.solid()
        right_shape.fill.fore_color.rgb = theme.accent_color
        right_text = right_shape.text_frame
        right_text.text = config.get('right_title', 'Option B')

        return slide

    def _create_timeline_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create timeline slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._apply_background(slide, theme, config.get('background'))

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = config.get('title', 'Timeline')
        title_para = title_frame.paragraphs[0]
        title_para.font.name = theme.title_font
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = theme.primary_color
        title_para.font.bold = True

        # Timeline events
        events = config.get('events', [])
        num_events = len(events)

        if num_events > 0:
            # Draw timeline line
            line = slide.shapes.add_connector(
                MSO_LINE.STRAIGHT_CONNECTOR_1,
                Inches(1), Inches(3.5),
                Inches(9), Inches(3.5)
            )
            line.line.color.rgb = theme.primary_color
            line.line.width = Pt(3)

            # Add event markers and labels
            spacing = 8.0 / (num_events - 1) if num_events > 1 else 0

            for i, event in enumerate(events):
                x_pos = 1 + (i * spacing)

                # Marker circle
                marker = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x_pos - 0.15), Inches(3.35),
                    Inches(0.3), Inches(0.3)
                )
                marker.fill.solid()
                marker.fill.fore_color.rgb = theme.accent_color
                marker.line.color.rgb = theme.primary_color

                # Event label
                label_box = slide.shapes.add_textbox(
                    Inches(x_pos - 0.5), Inches(4),
                    Inches(1), Inches(1)
                )
                label_frame = label_box.text_frame
                label_frame.text = event
                label_para = label_frame.paragraphs[0]
                label_para.font.size = Pt(10)
                label_para.alignment = PP_ALIGN.CENTER
                label_para.font.color.rgb = theme.text_color

        return slide

    def _create_image_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create slide with image"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._apply_background(slide, theme, config.get('background'))

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = config.get('title', 'Image Slide')
        title_para = title_frame.paragraphs[0]
        title_para.font.name = theme.title_font
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = theme.primary_color
        title_para.font.bold = True

        # Add image if path provided
        if 'image_path' in config and os.path.exists(config['image_path']):
            img_path = config['image_path']
            left = Inches(config.get('image_left', 2))
            top = Inches(config.get('image_top', 2))
            width = Inches(config.get('image_width', 6))

            pic = slide.shapes.add_picture(img_path, left, top, width=width)

        # Add caption if provided
        if 'caption' in config:
            caption_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6),
                Inches(9), Inches(0.5)
            )
            caption_frame = caption_box.text_frame
            caption_frame.text = config['caption']
            caption_para = caption_frame.paragraphs[0]
            caption_para.font.size = Pt(14)
            caption_para.alignment = PP_ALIGN.CENTER
            caption_para.font.color.rgb = theme.text_color
            caption_para.font.italic = True

        return slide

    def _create_blank_slide(self, prs: Presentation, config: Dict, theme: AdvancedTheme):
        """Create blank slide for full custom content"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._apply_background(slide, theme, config.get('background'))

        return slide

    def _apply_background(self, slide, theme: AdvancedTheme, custom_bg=None):
        """Apply background to slide"""
        if custom_bg:
            if isinstance(custom_bg, tuple) and len(custom_bg) == 3:
                # Custom RGB color
                self._apply_solid_background(slide, RGBColor(*custom_bg))
            elif custom_bg == 'gradient' or theme.has_gradient:
                # Gradient background (requires XML manipulation)
                pass
        else:
            # Apply theme background
            self._apply_solid_background(slide, theme.background_color)

    def _apply_solid_background(self, slide, color: RGBColor):
        """Apply solid color background"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_transition(self, slide, transition_type: str, speed: str = 'medium'):
        """Add transition to slide (requires XML manipulation)"""
        # Transition implementation through XML
        # This is a simplified version - full implementation would require more XML work
        pass

    def _apply_theme_to_presentation(self, prs: Presentation, theme: AdvancedTheme):
        """Apply theme to entire presentation"""
        for slide in prs.slides:
            self._apply_background(slide, theme, None)
            # Update text colors throughout
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.font.color.type == MSO_THEME_COLOR.TEXT_1:
                            paragraph.font.color.rgb = theme.text_color

    def _update_slide_content(self, slide, new_content: Dict):
        """Update content of existing slide"""
        # Update title if provided
        if 'title' in new_content:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    # Assume first text shape is title
                    shape.text_frame.text = new_content['title']
                    break

        # Update bullets if provided
        if 'bullets' in new_content:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    # Clear existing content
                    text_frame.clear()
                    # Add new bullets
                    for i, bullet in enumerate(new_content['bullets']):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = bullet

    def _reorder_slides(self, prs: Presentation, new_order: List[int]):
        """Reorder slides according to new_order list"""
        # Create new slide ID list based on new order
        sldIdLst = prs.slides._sldIdLst
        new_sldIdLst = [sldIdLst[i] for i in new_order if i < len(sldIdLst)]

        # Replace the slide ID list
        for i, sldId in enumerate(new_sldIdLst):
            sldIdLst.insert(i, sldId)
            if i + len(new_sldIdLst) < len(sldIdLst):
                sldIdLst.remove(sldIdLst[i + len(new_sldIdLst)])

    def get_available_themes(self) -> List[str]:
        """Return list of available theme names"""
        return list(self.themes.keys())

    def get_theme_preview(self, theme_name: str) -> Dict:
        """Get theme color preview"""
        theme = self.themes.get(theme_name)
        if not theme:
            return {}

        return {
            'name': theme.name,
            'colors': {
                'primary': (theme.primary_color.r, theme.primary_color.g, theme.primary_color.b),
                'secondary': (theme.secondary_color.r, theme.secondary_color.g, theme.secondary_color.b),
                'accent': (theme.accent_color.r, theme.accent_color.g, theme.accent_color.b),
                'text': (theme.text_color.r, theme.text_color.g, theme.text_color.b),
                'background': (theme.background_color.r, theme.background_color.g, theme.background_color.b)
            },
            'fonts': {
                'title': theme.title_font,
                'body': theme.body_font
            }
        }

def main():
    """Command line interface for enhanced presentation generation"""
    if len(sys.argv) < 2:
        print("Enhanced PPT Creator - Full Customization Control")
        print("\nUsage:")
        print("  python ppt_creator_enhanced.py config.json")
        print("\nAvailable themes:")
        creator = EnhancedPPTCreator()
        for theme in creator.get_available_themes():
            print(f"  - {theme}")
        return

    config_file = sys.argv[1]

    try:
        with open(config_file, 'r') as f:
            config = json.load(f)

        creator = EnhancedPPTCreator()

        # Check if editing existing file
        if 'edit_file' in config:
            result_path = creator.edit_presentation(config['edit_file'], config)
            print(f"✅ Presentation edited successfully: {result_path}")
        else:
            result_path = creator.create_presentation(config)
            print(f"✅ Presentation created successfully: {result_path}")

        print(f"📊 Theme: {config.get('theme', 'corporate_blue')}")
        print(f"📋 Slides: {len(config.get('slides', []))}")

    except Exception as e:
        print(f"❌ Error: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()
